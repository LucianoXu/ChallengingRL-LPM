#!/usr/bin/env python3
"""Generate kickoff_slides.pptx — 4 slides with rendered formulas + env screenshots.

Regenerate after content edits with:
    python3 render_assets.py        # rebuilds PNGs in assets/
    python3 make_slides.py          # rebuilds the pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE   = Path(__file__).parent
ASSETS = HERE / "assets"
OUT    = HERE / "kickoff_slides.pptx"

# ── Palette ────────────────────────────────────────────────────────────
INK    = RGBColor(0x1A, 0x1A, 0x1A)
MUTED  = RGBColor(0x60, 0x60, 0x60)
ACCENT = RGBColor(0x1E, 0x4E, 0x8C)   # RUB-ish navy
RED    = RGBColor(0xB7, 0x2B, 0x2B)
BG_BOX = RGBColor(0xF2, 0xF2, 0xF4)
RULE   = RGBColor(0xCC, 0xCC, 0xCC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────
def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return tf


def para(tf, text, *, size=14, bold=False, italic=False, color=INK,
         name="Calibri", align=PP_ALIGN.LEFT, level=0, first=False, space_after=4):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    f = p.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = name
    return p


def rect(slide, left, top, width, height, fill=BG_BOX):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def header(slide, title, subtitle=None):
    tf = textbox(slide, 0.5, 0.3, 12.3, 0.85)
    para(tf, title, size=30, bold=True, color=ACCENT, first=True, space_after=2)
    if subtitle:
        para(tf, subtitle, size=16, italic=True, color=MUTED, space_after=0)
    rule = slide.shapes.add_connector(
        1, Inches(0.5), Inches(1.25), Inches(12.83), Inches(1.25)
    )
    rule.line.color.rgb = RULE
    rule.line.width = Pt(0.75)


def img(slide, name, left, top, width=None, height=None):
    """Embed a picture from ./assets, scaling proportionally if one dim given."""
    kw = {}
    if width is not None:  kw["width"]  = Inches(width)
    if height is not None: kw["height"] = Inches(height)
    return slide.shapes.add_picture(
        str(ASSETS / name), Inches(left), Inches(top), **kw
    )


# ═══════════════════════════════════════════════════════════════════════
# Slide 1 — Topic & motivation
# ═══════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
header(s1, "Exploration vs. Exploitation",
       "When curiosity gets stuck on noise")

# Lead paragraph
tf = textbox(s1, 0.5, 1.4, 12.3, 1.4)
para(tf, "Lab topic  (1 of 5):  ",
     size=15, bold=True, color=INK, first=True)
para(tf, "the fundamental RL trade-off — try new actions vs. exploit the current best policy.",
     size=15, color=INK, space_after=8)
para(tf,
     "The sparse-reward bottleneck.  When extrinsic rewards are rare, naive exploration fails. "
     "The classical fix is intrinsic motivation — give the agent a self-generated bonus "
     "for visiting novel states.",
     size=14, color=INK)

# Table of exploration families
rows, cols = 4, 3
tbl_shape = s1.shapes.add_table(
    rows, cols, Inches(0.5), Inches(3.05), Inches(12.3), Inches(1.95)
)
tbl = tbl_shape.table
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(4.8)
tbl.columns[2].width = Inches(4.5)
headers = ["Family", "Intrinsic bonus", "Known issue"]
data = [
    ("ε-greedy / Softmax",      "none (random)",                       "Undirected — fails on long-horizon / sparse tasks"),
    ("UCB / Thompson Sampling", "uncertainty-driven",                  "Hard to scale to deep nets without approximation"),
    ("RND / ICM  (curiosity)",  "dynamics-model prediction error",     "Noisy-TV problem  ↓"),
]
for j, h in enumerate(headers):
    c = tbl.cell(0, j)
    c.text = h
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT
    p = c.text_frame.paragraphs[0]
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = RGBColor(255,255,255)
for i, row in enumerate(data, start=1):
    for j, txt in enumerate(row):
        c = tbl.cell(i, j); c.text = txt
        p = c.text_frame.paragraphs[0]
        p.font.size = Pt(12); p.font.color.rgb = INK
        if i == 3 and j == 2:
            p.font.bold = True; p.font.color.rgb = RED

# Punchline
tf = textbox(s1, 0.5, 5.35, 12.3, 1.9)
para(tf, "The noisy-TV failure mode.",
     size=16, bold=True, color=RED, first=True, space_after=6)
para(tf,
     "A curiosity-driven agent placed near an inherently random source (TV static, dice, sensor "
     "noise) becomes permanently fixated — prediction error never falls, the bonus never decays, "
     "the agent never moves on.  This is the failure mode our project will study.",
     size=14, color=INK)


# ═══════════════════════════════════════════════════════════════════════
# Slide 2 — LPM paper (formulas as rendered math images)
# ═══════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
header(s2, "Paper:  Learning Progress Monitoring (LPM)",
       "Hou, An, Du — UC Merced — ICLR 2026 — arXiv:2509.25438")

# Key-shift box (taller to give formulas breathing room)
rect(s2, 0.5, 1.45, 12.3, 3.0, fill=BG_BOX)

tf = textbox(s2, 0.7, 1.55, 11.9, 0.4)
para(tf, "Key conceptual shift — reward improvement, not prediction error:",
     size=15, bold=True, color=INK, first=True)

# Row 1 — Curiosity (label left, formula right)
tf = textbox(s2, 0.9, 2.15, 3.7, 0.5)
para(tf, "Curiosity (RND, ICM):",
     size=14, italic=True, color=MUTED, first=True)
img(s2, "formula_curiosity.png", left=4.4, top=2.05, width=7.8)

# Row 2 — LPM (label left, formula right, accented)
tf = textbox(s2, 0.9, 3.45, 3.7, 0.5)
para(tf, "LPM (this paper):",
     size=14, italic=True, bold=True, color=ACCENT, first=True)
img(s2, "formula_lpm.png", left=4.4, top=3.30, width=7.8)

# Why-it-works (below the box)
tf = textbox(s2, 0.5, 4.65, 12.3, 1.35)
para(tf, "Why it escapes noisy-TVs.",
     size=15, bold=True, color=INK, first=True, space_after=4)
para(tf,
     "Unlearnable transitions never get more predictable → bonus → 0 → agent ignores the noise "
     "source.  Learnable transitions keep improving → positive bonus → agent explores them.",
     size=13, color=INK, space_after=4)
para(tf,
     "The pathology is removed by changing what we reward, not by patching the symptoms.",
     size=13, italic=True, color=ACCENT)

# Theoretical guarantee
tf = textbox(s2, 0.5, 6.05, 12.3, 0.85)
para(tf, "Theoretical guarantee.",
     size=15, bold=True, color=INK, first=True, space_after=4)
para(tf,
     "LPM's bonus is zero-equivariant and a monotone indicator of Information Gain — "
     "matching principled Bayesian exploration, without posterior inference.",
     size=13, color=INK)

# Bottom metadata strip
tf = textbox(s2, 0.5, 6.95, 12.3, 0.45)
para(tf,
     "Architecture: dual-network (world model + error model)   ·   "
     "Code: github.com/Akuna23Matata/LPM_exploration   ·   "
     "Envs: Noisy-MNIST · MiniWorld 3D · Atari",
     size=11, color=MUTED, first=True)


# ═══════════════════════════════════════════════════════════════════════
# Slide 3 — NEW: Experimental environments
# ═══════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
header(s3, "Experimental Environments",
       "MiniGrid (our extension)   ◆   MiniWorld (paper's main env)")

# ── Left column: MiniGrid ──────────────────────────────────────────────
tf = textbox(s3, 0.5, 1.45, 6.0, 0.45)
para(tf, "MiniGrid  —  2D top-down gridworld",
     size=17, bold=True, color=ACCENT, first=True)

# Centered screenshot, height-driven sizing → 3×3 square
img(s3, "minigrid_doorkey.png", left=2.0, top=2.05, height=3.0)

tf = textbox(s3, 0.5, 5.2, 6.0, 0.35)
para(tf, "MiniGrid-DoorKey-8x8-v0",
     size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER, first=True)

# Bullets
tf = textbox(s3, 0.5, 5.7, 6.0, 1.7)
para(tf, "• Symbolic obs (or 7×7 partial RGB cone)",
     size=12, color=INK, first=True, space_after=3)
para(tf, "• Pure Python, no OpenGL — CPU only", size=12, color=INK, space_after=3)
para(tf, "• Sparse-reward, no noise → clean baseline for our Stage-2 extension",
     size=12, color=INK, space_after=3)
para(tf, "• Maintainer: Farama Foundation", size=11, italic=True, color=MUTED)

# ── Right column: MiniWorld ────────────────────────────────────────────
tf = textbox(s3, 6.83, 1.45, 6.0, 0.45)
para(tf, "MiniWorld  —  3D first-person maze",
     size=17, bold=True, color=ACCENT, first=True)

# Screenshot, height 3" → width auto = 4" (4:3 aspect)
img(s3, "miniworld_mazes3.png", left=7.83, top=2.05, height=3.0)

tf = textbox(s3, 6.83, 5.2, 6.0, 0.35)
para(tf, "MiniWorld-MazeS3-v0   (rendered 160×120 RGB in paper)",
     size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER, first=True)

tf = textbox(s3, 6.83, 5.7, 6.0, 1.7)
para(tf, "• Pixel obs (RGB) — first-person camera",
     size=12, color=INK, first=True, space_after=3)
para(tf, "• Pyglet + OpenGL — CPU + light GPU; CGL backend on Apple Silicon",
     size=12, color=INK, space_after=3)
para(tf, "• Paper's main env — supports noisy-wall / action-noise demos",
     size=12, color=INK, space_after=3)
para(tf, "• Maintainer: Farama Foundation", size=11, italic=True, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════
# Slide 4 — Project plan
# ═══════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
header(s4, "Project Plan",
       "Reproduce  →  controlled study")

# Stage 1
tf = textbox(s4, 0.5, 1.4, 12.3, 1.5)
para(tf, "Stage 1 — Reproduction   (≈ 2–3 weeks)",
     size=17, bold=True, color=ACCENT, first=True, space_after=6)
para(tf, "• Run LPM on Noisy-MNIST  (CPU-friendly toy)  and MiniWorld 3D maze.",
     size=13, color=INK, space_after=2)
para(tf, "• Match published curves within seed variance against ≥ 2 baselines (RND, vanilla DQN).",
     size=13, color=INK, space_after=2)
para(tf, "• Deliverable:  forked working repo  +  reproduction notebook with curves.",
     size=13, italic=True, color=MUTED)

# Stage 2
tf = textbox(s4, 0.5, 3.05, 12.3, 2.1)
para(tf, "Stage 2 — Controlled Study   (≈ 1–2 weeks)",
     size=17, bold=True, color=ACCENT, first=True, space_after=6)
para(tf, "1.  Noise-strength sweep on Noisy-MNIST  —  how does LPM's gain scale as the unlearnable "
         "fraction grows from 0 % to 100 %?", size=13, color=INK, space_after=2)
para(tf, "2.  Extension to MiniGrid  (not in paper)  —  does LPM still help on a cleaner sparse-reward "
         "env, or is its advantage noise-specific?", size=13, color=INK, space_after=2)
para(tf, "3.  Ablation — replace the error-model with a running average of prediction errors.  "
         "How much of LPM's gain comes from the dual-network design vs. the "
         "\"reward-improvement-not-error\" insight?", size=13, color=INK)

# Bottom: Deliverables (left) + Risk mitigation (right)
tf = textbox(s4, 0.5, 5.35, 6.0, 1.9)
para(tf, "Deliverables  &  Compute",
     size=15, bold=True, color=INK, first=True, space_after=6)
para(tf, "• Forked & modified code repo", size=13, color=INK, space_after=2)
para(tf, "• Written report",              size=13, color=INK, space_after=2)
para(tf, "• Group presentation",          size=13, color=INK, space_after=6)
para(tf, "All primary experiments target single-GPU laptop / CPU — no cluster.",
     size=12, italic=True, color=MUTED)

tf = textbox(s4, 6.83, 5.35, 6.0, 1.9)
para(tf, "Risk Mitigation",
     size=15, bold=True, color=RED, first=True, space_after=6)
para(tf,
     "Backup paper:  ETD — Episodic Novelty Through Temporal Distance "
     "(Liu et al., ICLR 2025, arXiv:2501.15418).",
     size=13, color=INK, space_after=4)
para(tf,
     "Same noisy-TV motivation via temporal-distance contrastive learning.  "
     "Drop-in replacement if LPM code turns out unrunnable.",
     size=12, italic=True, color=MUTED)


# ── Save ───────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Wrote {OUT}")
